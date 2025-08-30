"""Export services for generating reports in various formats."""

import io
import uuid
from datetime import datetime
from typing import Dict, List, Any, Optional, BinaryIO
from pathlib import Path

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import LineChart, Reference
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.core.config import settings
from app.core.observability import trace_function


class ExcelExporter:
    """Excel export functionality."""
    
    def __init__(self) -> None:
        """Initialize Excel exporter."""
        pass
    
    @trace_function("excel_exporter.create_financial_report")
    def create_financial_report(
        self,
        company_data: Dict[str, Any],
        financial_data: pd.DataFrame,
        valuation_data: Optional[Dict[str, Any]] = None,
        charts_data: Optional[List[Dict[str, Any]]] = None,
    ) -> bytes:
        """Create comprehensive financial report in Excel."""
        # Create workbook
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Create sheets
        self._create_summary_sheet(wb, company_data, valuation_data)
        self._create_financials_sheet(wb, financial_data)
        
        if valuation_data:
            self._create_valuation_sheet(wb, valuation_data)
        
        if charts_data:
            self._create_charts_sheet(wb, charts_data)
        
        # Save to bytes
        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        
        return buffer.getvalue()
    
    def _create_summary_sheet(
        self, 
        wb: Workbook, 
        company_data: Dict[str, Any],
        valuation_data: Optional[Dict[str, Any]]
    ) -> None:
        """Create executive summary sheet."""
        ws = wb.create_sheet("Executive Summary", 0)
        
        # Header
        ws['A1'] = f"Financial Analysis - {company_data.get('name', 'Company')}"
        ws['A1'].font = Font(size=16, bold=True)
        ws['A2'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Company info
        row = 4
        ws[f'A{row}'] = "Company Information"
        ws[f'A{row}'].font = Font(bold=True)
        
        info_items = [
            ("Symbol", company_data.get('symbol', 'N/A')),
            ("Sector", company_data.get('sector', 'N/A')),
            ("Market Cap", f"${company_data.get('market_cap', 0):,.0f}"),
            ("Enterprise Value", f"${company_data.get('enterprise_value', 0):,.0f}"),
        ]
        
        for i, (label, value) in enumerate(info_items):
            ws[f'A{row + 1 + i}'] = label
            ws[f'B{row + 1 + i}'] = value
        
        # Valuation summary
        if valuation_data:
            row += len(info_items) + 3
            ws[f'A{row}'] = "Valuation Summary"
            ws[f'A{row}'].font = Font(bold=True)
            
            val_items = [
                ("Target Price", f"${valuation_data.get('share_price', 0):.2f}"),
                ("Current Price", f"${valuation_data.get('current_price', 0):.2f}"),
                ("Upside/Downside", f"{valuation_data.get('upside', 0):.1f}%"),
                ("Model", valuation_data.get('model_type', 'N/A')),
            ]
            
            for i, (label, value) in enumerate(val_items):
                ws[f'A{row + 1 + i}'] = label
                ws[f'B{row + 1 + i}'] = value
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
    
    def _create_financials_sheet(self, wb: Workbook, financial_data: pd.DataFrame) -> None:
        """Create financial data sheet."""
        ws = wb.create_sheet("Financial Data")
        
        # Headers
        ws['A1'] = "Financial Data"
        ws['A1'].font = Font(size=14, bold=True)
        
        # Write DataFrame to sheet
        for r_idx, row in enumerate(financial_data.itertuples(), 3):
            for c_idx, value in enumerate(row[1:], 1):  # Skip index
                ws.cell(row=r_idx, column=c_idx, value=value)
        
        # Write column headers
        for c_idx, col_name in enumerate(financial_data.columns, 1):
            cell = ws.cell(row=2, column=c_idx, value=col_name)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
    
    def _create_valuation_sheet(self, wb: Workbook, valuation_data: Dict[str, Any]) -> None:
        """Create valuation analysis sheet."""
        ws = wb.create_sheet("Valuation Analysis")
        
        ws['A1'] = "DCF Valuation Analysis"
        ws['A1'].font = Font(size=14, bold=True)
        
        # Key metrics
        row = 3
        metrics = [
            ("Enterprise Value", f"${valuation_data.get('enterprise_value', 0):,.0f}"),
            ("Equity Value", f"${valuation_data.get('equity_value', 0):,.0f}"),
            ("Share Price", f"${valuation_data.get('share_price', 0):.2f}"),
            ("Terminal Value", f"${valuation_data.get('terminal_value', 0):,.0f}"),
            ("PV of Explicit Period", f"${valuation_data.get('pv_explicit_period', 0):,.0f}"),
            ("PV of Terminal Value", f"${valuation_data.get('pv_terminal_value', 0):,.0f}"),
        ]
        
        for i, (label, value) in enumerate(metrics):
            ws[f'A{row + i}'] = label
            ws[f'B{row + i}'] = value
        
        # Assumptions
        if 'assumptions' in valuation_data:
            row += len(metrics) + 2
            ws[f'A{row}'] = "Key Assumptions"
            ws[f'A{row}'].font = Font(bold=True)
            
            assumptions = valuation_data['assumptions']
            for i, (key, value) in enumerate(assumptions.items()):
                ws[f'A{row + 1 + i}'] = key.replace('_', ' ').title()
                if isinstance(value, float) and 0 < value < 1:
                    ws[f'B{row + 1 + i}'] = f"{value:.1%}"
                else:
                    ws[f'B{row + 1 + i}'] = str(value)
    
    def _create_charts_sheet(self, wb: Workbook, charts_data: List[Dict[str, Any]]) -> None:
        """Create charts sheet."""
        ws = wb.create_sheet("Charts")
        
        ws['A1'] = "Financial Charts"
        ws['A1'].font = Font(size=14, bold=True)
        
        # For now, just add chart descriptions
        # In production, would create actual Excel charts
        row = 3
        for i, chart in enumerate(charts_data):
            ws[f'A{row + i}'] = f"Chart {i+1}: {chart.get('title', 'Untitled')}"


class PowerPointExporter:
    """PowerPoint export functionality."""
    
    def __init__(self) -> None:
        """Initialize PowerPoint exporter."""
        pass
    
    @trace_function("pptx_exporter.create_investment_deck")
    def create_investment_deck(
        self,
        company_data: Dict[str, Any],
        valuation_data: Dict[str, Any],
        financial_data: pd.DataFrame,
        charts_data: Optional[List[Dict[str, Any]]] = None,
    ) -> bytes:
        """Create investment thesis PowerPoint deck."""
        prs = Presentation()
        
        # Title slide
        self._create_title_slide(prs, company_data)
        
        # Executive summary
        self._create_executive_summary_slide(prs, company_data, valuation_data)
        
        # Financial overview
        self._create_financial_overview_slide(prs, financial_data)
        
        # Valuation slide
        self._create_valuation_slide(prs, valuation_data)
        
        # Charts slides
        if charts_data:
            for chart in charts_data:
                self._create_chart_slide(prs, chart)
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return buffer.getvalue()
    
    def _create_title_slide(self, prs: Presentation, company_data: Dict[str, Any]) -> None:
        """Create title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Investment Analysis: {company_data.get('name', 'Company')}"
        subtitle.text = f"Symbol: {company_data.get('symbol', 'N/A')} | {datetime.now().strftime('%B %Y')}"
    
    def _create_executive_summary_slide(
        self, 
        prs: Presentation, 
        company_data: Dict[str, Any],
        valuation_data: Dict[str, Any]
    ) -> None:
        """Create executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Key points
        points = [
            f"Target Price: ${valuation_data.get('share_price', 0):.2f}",
            f"Current Price: ${valuation_data.get('current_price', 0):.2f}",
            f"Upside Potential: {valuation_data.get('upside', 0):.1f}%",
            f"Market Cap: ${company_data.get('market_cap', 0):,.0f}",
            f"Valuation Method: {valuation_data.get('model_type', 'DCF')}",
        ]
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    
    def _create_financial_overview_slide(self, prs: Presentation, financial_data: pd.DataFrame) -> None:
        """Create financial overview slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Financial Overview"
        
        # Add table with key metrics
        if not financial_data.empty:
            content = slide.placeholders[1]
            tf = content.text_frame
            
            # Show recent metrics
            recent_data = financial_data.tail(3) if len(financial_data) > 3 else financial_data
            
            for _, row in recent_data.iterrows():
                p = tf.add_paragraph()
                p.text = f"Year {row.get('year', 'N/A')}: Revenue ${row.get('revenue', 0):,.0f}"
                p.level = 0
    
    def _create_valuation_slide(self, prs: Presentation, valuation_data: Dict[str, Any]) -> None:
        """Create valuation analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Valuation Analysis"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Valuation metrics
        metrics = [
            f"Enterprise Value: ${valuation_data.get('enterprise_value', 0):,.0f}",
            f"Equity Value: ${valuation_data.get('equity_value', 0):,.0f}",
            f"Target Share Price: ${valuation_data.get('share_price', 0):.2f}",
            f"Terminal Value: ${valuation_data.get('terminal_value', 0):,.0f}",
        ]
        
        for metric in metrics:
            p = tf.add_paragraph()
            p.text = metric
            p.level = 0
    
    def _create_chart_slide(self, prs: Presentation, chart_data: Dict[str, Any]) -> None:
        """Create chart slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = chart_data.get('title', 'Chart')
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = chart_data.get('description', 'Chart description would appear here.')


class PDFExporter:
    """PDF export functionality."""
    
    def __init__(self) -> None:
        """Initialize PDF exporter."""
        self.styles = getSampleStyleSheet()
    
    @trace_function("pdf_exporter.create_research_report")
    def create_research_report(
        self,
        company_data: Dict[str, Any],
        analysis_text: str,
        financial_data: pd.DataFrame,
        valuation_data: Optional[Dict[str, Any]] = None,
        citations: Optional[List[Dict[str, Any]]] = None,
    ) -> bytes:
        """Create research report in PDF format."""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        
        story = []
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1,  # Center alignment
        )
        
        title = Paragraph(
            f"Financial Analysis Report<br/>{company_data.get('name', 'Company')} ({company_data.get('symbol', 'N/A')})",
            title_style
        )
        story.append(title)
        story.append(Spacer(1, 12))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", self.styles['Heading2']))
        
        if valuation_data:
            summary_text = f"""
            Target Price: ${valuation_data.get('share_price', 0):.2f}<br/>
            Current Price: ${valuation_data.get('current_price', 0):.2f}<br/>
            Upside Potential: {valuation_data.get('upside', 0):.1f}%<br/>
            Valuation Method: {valuation_data.get('model_type', 'DCF')}
            """
            story.append(Paragraph(summary_text, self.styles['Normal']))
        
        story.append(Spacer(1, 12))
        
        # Analysis
        story.append(Paragraph("Analysis", self.styles['Heading2']))
        story.append(Paragraph(analysis_text, self.styles['Normal']))
        story.append(Spacer(1, 12))
        
        # Financial Data Table
        if not financial_data.empty:
            story.append(Paragraph("Financial Data", self.styles['Heading2']))
            
            # Create table data
            table_data = [financial_data.columns.tolist()]
            for _, row in financial_data.head(10).iterrows():  # Limit to 10 rows
                table_data.append([str(val) for val in row.tolist()])
            
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(table)
            story.append(Spacer(1, 12))
        
        # Citations
        if citations:
            story.append(Paragraph("Sources", self.styles['Heading2']))
            for i, citation in enumerate(citations, 1):
                citation_text = f"{i}. {citation.get('source', 'Unknown')} - {citation.get('locator', '')}"
                story.append(Paragraph(citation_text, self.styles['Normal']))
        
        # Footer
        story.append(Spacer(1, 24))
        footer_text = f"Generated by AI Financial Analyst on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        story.append(Paragraph(footer_text, self.styles['Normal']))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        return buffer.getvalue()


class ExportService:
    """Main export service coordinating different exporters."""
    
    def __init__(self) -> None:
        """Initialize export service."""
        self.excel_exporter = ExcelExporter()
        self.pptx_exporter = PowerPointExporter()
        self.pdf_exporter = PDFExporter()
    
    @trace_function("export_service.export_analysis")
    async def export_analysis(
        self,
        format_type: str,  # excel, pptx, pdf
        data: Dict[str, Any],
    ) -> bytes:
        """Export analysis in specified format."""
        if format_type.lower() == 'excel':
            return self.excel_exporter.create_financial_report(
                company_data=data.get('company', {}),
                financial_data=pd.DataFrame(data.get('financial_data', [])),
                valuation_data=data.get('valuation', {}),
                charts_data=data.get('charts', []),
            )
        
        elif format_type.lower() == 'pptx':
            return self.pptx_exporter.create_investment_deck(
                company_data=data.get('company', {}),
                valuation_data=data.get('valuation', {}),
                financial_data=pd.DataFrame(data.get('financial_data', [])),
                charts_data=data.get('charts', []),
            )
        
        elif format_type.lower() == 'pdf':
            return self.pdf_exporter.create_research_report(
                company_data=data.get('company', {}),
                analysis_text=data.get('analysis_text', ''),
                financial_data=pd.DataFrame(data.get('financial_data', [])),
                valuation_data=data.get('valuation', {}),
                citations=data.get('citations', []),
            )
        
        else:
            raise ValueError(f"Unsupported export format: {format_type}")
    
    def get_filename(self, format_type: str, company_symbol: str = "analysis") -> str:
        """Generate filename for export."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"{company_symbol}_analysis_{timestamp}.{format_type.lower()}"
